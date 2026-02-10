from pptx import Presentation
from pptx.util import Inches
import os
import glob

# Define paths
OUTPUT_PATH = "Paper/presentation.pptx"
SLIDES_DIR = "Paper/slides"

def create_presentation_from_images():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio standard for most modern screens
    # (PowerPoint defaults to 16:9, but let's confirm width/height if needed, default is usually fine)
    # Default is 10x7.5 inches (4:3). Metropolis/Beamer usually outputs 4:3 or 16:9 depending on settings.
    # We should check the aspect ratio of the images or just fit them.
    # Beamer default is 128mm x 96mm (4:3).
    # Since we are pasting images, we want the slide to match the image AR.
    # Let's assume standard Beamer 4:3 for now, but if Metropolis is 16:9 (often is), we might need to adjust.
    # Actually, pdftoppm preserves the aspect ratio.
    # Let's set the PPTX slide size to the image size of the first image to be pixel perfect.
    
    image_files = sorted(glob.glob(os.path.join(SLIDES_DIR, "*.png")))
    
    if not image_files:
        print("No images found!")
        return

    # Create a dummy presentation to get default slide width/height
    # Then we will adjust it based on the first image.
    # python-pptx allows setting slide_width and slide_height.
    
    # Let's blindly add images and maximize them.
    
    # Use blank layout
    BLANK_LAYOUT_INDEX = 6 
    
    for img_path in image_files:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
        
        # We want to fill the slide.
        # Ideally we read image size but python-pptx doesn't give us image dims easily without PIL.
        # Since we installed Pillow earlier implicitly (deps of python-pptx), we can use it.
        try:
            from PIL import Image
            with Image.open(img_path) as im:
                 width_px, height_px = im.size
                 
            # Convert pixels to EMUs (English Metric Units) or Inches for PPTX
            # PPTX uses EMUs. 914400 EMUs per inch.
            # Resolution of pdftoppm was 300 DPI.
            
            width_inch = width_px / 300.0
            height_inch = height_px / 300.0
            
            prs.slide_width = int(width_inch * 914400)
            prs.slide_height = int(height_inch * 914400)
            
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        except Exception as e:
            print(f"Error processing {img_path}: {e}")
            # Fallback if PIL fails
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width)

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to {OUTPUT_PATH}")

if __name__ == "__main__":
    create_presentation_from_images()
