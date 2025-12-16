from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Define paths
OUTPUT_PATH = "Paper/presentation.pptx"
IMG_DIR = "results/graphs"

def insert_image_into_placeholder(slide, placeholder_idx, img_path):
    placeholder = slide.placeholders[placeholder_idx]
    # Calculate aspect ratio preserving dimensions
    im = slide.shapes.add_picture(img_path, placeholder.left, placeholder.top, width=placeholder.width)
    # Optional: Center vertically/horizontally if aspect ratio differs (omitted for simplicity, just width-fit)

def create_presentation():
    prs = Presentation()
    
    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Enhancing QoS in Dense IEEE 802.11ax Settings\nusing a Dynamic Airtime-Based Soft Admission Control Mechanism"
    subtitle.text = "Dayanand Ambawade & Rohan Pawar\n\n10th International Conference on Systems, Control and Communications (ICSCC 2025)\nNagoya University, Japan\nDecember 2025"

    # --- Slide 2: Outline ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Outline"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Introduction & Motivation"
    p = tf.add_paragraph()
    p.text = "System Architecture"
    p = tf.add_paragraph()
    p.text = "Proposed AS-CAC Framework"
    p = tf.add_paragraph()
    p.text = "Performance Evaluation"
    p = tf.add_paragraph()
    p.text = "Conclusion"

    # --- Slide 3: Motivation ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Motivation"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "The Challenge"
    p = tf.paragraphs[0]
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "IEEE 802.11ax (Wi-Fi 6) improves spectral efficiency but struggles under saturation in dense environments."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "QoS Issues"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Real-time applications (VoIP, Video) are sensitive to high delay and jitter caused by congestion."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Limitation of Existing CAC"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Traditional schemes ignore traffic heterogeneity (IoT vs 4K Video) and varying MCS."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Research Goal"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Develop admission control to maximize airtime utilization while guaranteeing strict QoS."
    p.level = 1

    # --- Slide 4: System Architecture (Text + Image) ---
    slide_layout = prs.slide_layouts[3] # Two Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Research Test Bed & Components"
    
    # Left: Text
    tf = slide.placeholders[1].text_frame
    tf.text = "Simulation Components (ns-3):"
    p = tf.add_paragraph()
    p.text = "AP Node: Wi-Fi 6 (802.11ax), 80 MHz, 5 GHz"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Stations: 25-50 users in dense grid"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Traffic Generators:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "VoIP (UDP)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Video (UDP)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Bursty/Web (TCP/UDP)"
    p.level = 2
    
    # Right: Image
    img_path = os.path.join(IMG_DIR, "ns3_network_topology.png")
    insert_image_into_placeholder(slide, 2, img_path)

    # --- Slide 5: System Model (Text + Image) ---
    slide_layout = prs.slide_layouts[3] # Two Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Model and Traffic"
    
    # Left: Text
    tf = slide.placeholders[1].text_frame
    tf.text = "Metric: Airtime Utilization"
    p = tf.add_paragraph()
    p.text = "Calculated based on PHY rate and overhead."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Traffic Mix:"
    
    p = tf.add_paragraph()
    p.text = "VoIP: High Priority, Low Bandwidth"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Video: Medium Priority, High Bandwidth"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Best Effort: Low Priority, Bursty"
    p.level = 1
    
    # Right: Image
    img_path = os.path.join(IMG_DIR, "ap_station_distribution.png")
    insert_image_into_placeholder(slide, 2, img_path)

    # --- Slide 6: Proposed Solution: Soft CAC ---
    slide_layout = prs.slide_layouts[3] # Two Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Proposed Solution: Soft CAC"
    
    # Left: Table-like Text
    tf = slide.placeholders[1].text_frame
    tf.text = "Concept: Priority-Aware Thresholds"
    
    p = tf.add_paragraph()
    p.text = "Instead of a single hard limit (e.g., 80%):"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "VoIP (High): 90% Threshold"
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Video (Medium): 80% Threshold"
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Best Effort (Low): 95% Threshold"
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Benefit: Best-effort traffic fills gaps without blocking high-priority flows."
    p.level = 1

    # Right: Image
    img_path = os.path.join(IMG_DIR, "soft_vs_hard_cac_comparison.png")
    insert_image_into_placeholder(slide, 2, img_path)

    # --- Slide 7: AS-CAC+ Adaptive ---
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Algorithm: AS-CAC+ (Adaptive)"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Dynamic Threshold Adjustment"
    p = tf.add_paragraph()
    p.text = "Monitors Packet Error Rate (PER) and Utilization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Adjusts Best-Effort threshold in real-time."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Logic:"
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "If PER > 5%: Decrease Threshold (Protect QoS)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "If PER < 2%: Increase Threshold (Utilize Capacity)"
    p.level = 2
    
    img_path = os.path.join(IMG_DIR, "ascac_threshold_evolution.png")
    insert_image_into_placeholder(slide, 2, img_path)

    # --- Slide 8: Latency Results ---
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Simulation Results: Latency Impact"
    
    img_path = os.path.join(IMG_DIR, "cac_vs_no_cac.png")
    # Centered large image
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(7)
    slide.shapes.add_picture(img_path, left, top, width=width)
    
    # Text box caption
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Without CAC, latency spikes > 45ms. AS-CAC maintains < 2ms."
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 9: Comprehensive Analysis ---
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Comprehensive Performance Analysis"
    
    img_path = os.path.join(IMG_DIR, "ascac_comprehensive_4panel.png")
    left = Inches(1)
    top = Inches(1.2)
    width = Inches(8) # Max width
    pic = slide.shapes.add_picture(img_path, left, top, width=width)

    # --- Slide 10: Multi-Dimensional Superiority ---
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Multi-Dimensional Superiority"

    # Left: Image
    img_path = os.path.join(IMG_DIR, "ascac_radar_chart.png")
    insert_image_into_placeholder(slide, 1, img_path)
    
    # Right: Text
    tf = slide.placeholders[2].text_frame
    tf.text = "Why AS-CAC+ Wins:"
    p = tf.add_paragraph()
    p.text = "Adaptability"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Reacts to interference instantly."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Utilization"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "97.4% vs 78% (Hard CAC)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Safety"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Keeps VoIP latency < 2ms."
    p.level = 1

    # --- Slide 11: Conclusion ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"

    tf = slide.placeholders[1].text_frame
    tf.text = "Summary"
    p = tf.add_paragraph()
    p.text = "AS-CAC+ transforms admission control from static to dynamic."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Safely unlocks 19.2% more capacity."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Key Achievements"
    p = tf.add_paragraph()
    p.text = "19.2% Throughput improvement."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "97.4% Channel Utilization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Strict QoS for VoIP validated."
    p.level = 1

    # --- Slide 12: Thank You ---
    slide_layout = prs.slide_layouts[0] # Title Slide (using as end slide)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Thank You!"
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Questions?\n\nDayanand Ambawade & Rohan Pawar\nSardar Patel Institute of Technology"

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to {OUTPUT_PATH}")

if __name__ == "__main__":
    create_presentation()
